#!/usr/bin/env python3
"""
DECP Platform Presentation Generator
Creates a comprehensive PowerPoint presentation for the DECP Platform project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define colors
DARK_BLUE = RGBColor(0x1a, 0x23, 0x7e)
LIGHT_BLUE = RGBColor(0x39, 0x49, 0xab)
ACCENT = RGBColor(0x00, 0x96, 0xc7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
ORANGE = RGBColor(0xFF, 0x98, 0x00)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = ACCENT
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullet_points):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)
        p.level = 0
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a slide with two columns"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    
    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    
    # Right column content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2), Inches(6), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
    
    return slide

# Slide 1: Title
add_title_slide(prs, 
    "Department Engagement & Career Platform (DECP)",
    "Connecting Students, Alumni & Faculty | University of Peradeniya")

# Slide 2: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "Students lack direct connection with alumni for mentorship",
    "No centralized platform for job opportunities within the university network",
    "Research collaboration is often limited by lack of visibility",
    "Department events have low attendance due to poor notification systems",
    "No single platform for academic and professional networking",
    "Existing platforms (LinkedIn, Facebook) lack academic-specific features"
])

# Slide 3: Solution Overview
add_content_slide(prs, "Solution: DECP Platform", [
    "Comprehensive social networking platform for academic communities",
    "Connects current students with alumni for mentorship and career guidance",
    "Centralized job and internship board with university network focus",
    "Research collaboration workspace with document sharing",
    "Event management system with automated reminders",
    "Real-time messaging for seamless communication",
    "Analytics dashboard for administrators"
])

# Slide 4: Key Features
add_two_column_slide(prs, "Key Features",
    "Social Features",
    [
        "User profiles with academic verification",
        "Social feed with posts, likes, comments",
        "Connection management",
        "Real-time messaging",
        "Media sharing (images/videos)",
        "Push notifications"
    ],
    "Professional Features",
    [
        "Job & internship board",
        "Application tracking system",
        "Research project collaboration",
        "Document sharing & version control",
        "Event management & RSVP",
        "Analytics dashboard"
    ]
)

# Slide 5: Architecture Overview
add_content_slide(prs, "System Architecture", [
    "Microservices Architecture with 10 independent services",
    "API Gateway for centralized routing and authentication",
    "Service-Oriented Architecture (SOA) for modularity",
    "Event-driven communication using RabbitMQ",
    "Real-time messaging with Socket.io",
    "Database per service pattern with PostgreSQL",
    "Redis for caching and session management",
    "Docker containers with Kubernetes orchestration"
])

# Slide 6: Technology Stack
add_two_column_slide(prs, "Technology Stack",
    "Backend Technologies",
    [
        "Node.js 20+ with Express.js",
        "TypeScript for type safety",
        "PostgreSQL 15+ (Primary DB)",
        "Redis (Caching & Sessions)",
        "RabbitMQ (Message Queue)",
        "Socket.io (Real-time)",
        "JWT & OAuth 2.0 (Auth)",
        "Sequelize ORM"
    ],
    "Frontend & DevOps",
    [
        "React 18+ with TypeScript",
        "React Native (Mobile)",
        "Redux Toolkit (State)",
        "Material-UI v5 + Tailwind",
        "Docker & Kubernetes",
        "AWS Cloud Infrastructure",
        "Terraform (IaC)",
        "GitHub Actions (CI/CD)"
    ]
)

# Slide 7: Microservices
add_content_slide(prs, "Microservices Architecture", [
    "API Gateway (Port: 3000) - Routing, Auth, Rate Limiting",
    "Auth Service (Port: 3001) - Registration, Login, OAuth",
    "User Service (Port: 3002) - Profiles, Connections",
    "Feed Service (Port: 3003) - Posts, Media, Engagement",
    "Jobs Service (Port: 3004) - Job postings, Applications",
    "Events Service (Port: 3005) - Event management, RSVP",
    "Research Service (Port: 3006) - Projects, Documents",
    "Messaging Service (Port: 3007) - Chat, Real-time",
    "Notification Service (Port: 3008) - Push, Email alerts",
    "Analytics Service (Port: 3009) - Metrics, Dashboards"
])

# Slide 8: Quality Attributes
add_content_slide(prs, "Quality Attributes", [
    "Scalability: Horizontal scaling with Kubernetes HPA",
    "Availability: 99.9% uptime with multi-AZ deployment",
    "Performance: < 200ms API response time (p95)",
    "Security: OWASP Top 10 compliant, JWT authentication",
    "Maintainability: Clean code, comprehensive documentation",
    "Observability: Prometheus metrics, Grafana dashboards",
    "Testability: 85%+ test coverage across all services",
    "Usability: Responsive UI, intuitive UX"
])

# Slide 9: Security Architecture
add_content_slide(prs, "Security Implementation", [
    "JWT-based authentication with short-lived access tokens",
    "Role-Based Access Control (RBAC) for authorization",
    "Password hashing with bcrypt (12 rounds)",
    "HTTPS/TLS 1.3 for all communications",
    "Input validation and SQL injection prevention",
    "XSS and CSRF protection",
    "Rate limiting: 100 req/min per IP",
    "AWS WAF for DDoS protection",
    "Regular security scans and dependency updates"
])

# Slide 10: Research Findings
add_content_slide(prs, "Research & Analysis", [
    "Analyzed Facebook and LinkedIn architectures",
    "Identified gaps: Academic verification, Research tools",
    "Proposed improvements: Mentorship matching, Academic calendar",
    "Architecture decisions informed by industry best practices",
    "Adopted microservices pattern from Netflix, Amazon",
    "Used PostgreSQL similar to LinkedIn's Espresso",
    "Implemented real-time features inspired by Facebook",
    "Privacy-first design addressing data concerns"
])

# Slide 11: Demo Features
add_content_slide(prs, "Live Demo Features", [
    "User Registration & Login with email verification",
    "Create and interact with posts (like, comment, share)",
    "Browse and apply for job opportunities",
    "Create and RSVP to department events",
    "Research project collaboration workspace",
    "Real-time messaging between users",
    "Push notifications for important updates",
    "Analytics dashboard for administrators",
    "Mobile app with all features"
])

# Slide 12: Testing & Quality
add_content_slide(prs, "Testing & Quality Assurance", [
    "Unit Tests: Jest for all microservices (85%+ coverage)",
    "Integration Tests: Service communication verification",
    "E2E Tests: Cypress for web, Detox for mobile",
    "Performance Tests: k6 for load, stress, and spike testing",
    "Security Tests: OWASP Top 10 compliance verification",
    "CI/CD: Automated testing on every commit",
    "Code Quality: ESLint, Prettier, TypeScript strict mode",
    "Documentation: Comprehensive API docs and guides"
])

# Slide 13: Deployment
add_content_slide(prs, "Cloud Deployment", [
    "AWS Infrastructure: VPC, EKS, RDS, ElastiCache, S3",
    "Kubernetes: Container orchestration with auto-scaling",
    "Terraform: Infrastructure as Code for reproducibility",
    "CI/CD Pipeline: GitHub Actions for automated deployment",
    "Monitoring: Prometheus + Grafana for observability",
    "Logging: Centralized logging with Loki",
    "CDN: CloudFront for static assets and media",
    "SSL/TLS: Automated certificate management with cert-manager"
])

# Slide 14: Project Structure
add_content_slide(prs, "Project Deliverables", [
    "Complete microservices backend (10 services)",
    "React web application with responsive design",
    "React Native mobile application",
    "Docker containers and Kubernetes manifests",
    "Terraform scripts for AWS infrastructure",
    "Comprehensive test suite (Unit, Integration, E2E)",
    "Architecture diagrams (C4 Model, SOA, Deployment)",
    "Complete documentation (API, Development, Deployment)",
    "CI/CD pipelines and DevOps automation"
])

# Slide 15: Team & Roles
add_content_slide(prs, "Team Structure", [
    "Enterprise Architect: High-level system design, integration patterns",
    "Solution Architect: Technology stack, cloud infrastructure decisions",
    "Application Architect: API design, microservices communication",
    "Security Architect: Authentication, authorization, data protection",
    "DevOps Architect: CI/CD, deployment, monitoring, infrastructure"
])

# Slide 16: Future Enhancements
add_content_slide(prs, "Future Enhancements", [
    "AI-powered job recommendations using machine learning",
    "Video conferencing integration for virtual events",
    "Mobile offline support for better accessibility",
    "Advanced analytics with predictive insights",
    "LMS integration (Moodle, Canvas) for course data",
    "Blockchain-based certificate verification",
    "Multi-language support for international students",
    "Alumni donation and fundraising module"
])

# Slide 17: Conclusion
add_title_slide(prs,
    "Thank You!",
    "DECP Platform - Production Ready | Questions?")

# Save presentation
prs.save('E:\\Academic\\Semester_7\\CO528\\Mini Project\\DECP_Platform_Presentation.pptx')
print("[OK] Presentation created successfully!")
print("Location: E:\\Academic\\Semester_7\\CO528\\Mini Project\\DECP_Platform_Presentation.pptx")
print("\nSlides created:")
print("  1. Title Slide")
print("  2. Problem Statement")
print("  3. Solution Overview")
print("  4. Key Features")
print("  5. Architecture Overview")
print("  6. Technology Stack")
print("  7. Microservices")
print("  8. Quality Attributes")
print("  9. Security Architecture")
print(" 10. Research Findings")
print(" 11. Demo Features")
print(" 12. Testing & Quality")
print(" 13. Deployment")
print(" 14. Project Deliverables")
print(" 15. Team Structure")
print(" 16. Future Enhancements")
print(" 17. Thank You")
